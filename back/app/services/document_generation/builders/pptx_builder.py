import os
import logging
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml

from app.core.files_config import BRAND_CONFIG
from app.domain.entities.document_content import DocumentContent
from app.domain.interfaces.document_builder import IDocumentBuilder
from app.services.document_generation.template_engine import TemplateEngine

logger = logging.getLogger(__name__)


def add_textbox_text(slide, left, top, width, height, text, font_size_pt, font_color_rgb, bold=False, align=PP_ALIGN.LEFT, font_name="Arial", italic=False):
    """Inserta un cuadro de texto con formato de una sola línea o bloque de texto básico."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size_pt)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = font_color_rgb
    return txBox


class PptxBuilder(IDocumentBuilder):

    def __init__(self, engine: TemplateEngine):
        self._engine = engine
        # Inicializar colores corporativos por defecto
        self._primary_rgb = RGBColor(1, 30, 35)      # #011E23
        self._secondary_rgb = RGBColor(16, 71, 63)   # #10473F
        self._accent_rgb = RGBColor(26, 235, 159)    # #1AEB9F
        self._text_rgb = RGBColor(45, 55, 72)        # #2D3748
        self._zebra_rgb = RGBColor(244, 246, 246)    # #F4F6F6

    @property
    def output_format(self) -> str:
        return "pptx"

    def build(self, content: DocumentContent) -> bytes:
        brand = content.brand or "convertia"
        brand_cfg = BRAND_CONFIG.get(brand, BRAND_CONFIG["convertia"])
        template_path = brand_cfg["templates"]["presentation"]

        # Configurar colores desde la configuración de la marca si están disponibles
        brand_colors = brand_cfg.get("colors", {})
        if "primary" in brand_colors:
            self._primary_rgb = self._hex_to_rgb(brand_colors["primary"])
        if "secondary_dark_green" in brand_colors:
            self._secondary_rgb = self._hex_to_rgb(brand_colors["secondary_dark_green"])
        if "accent_green" in brand_colors:
            self._accent_rgb = self._hex_to_rgb(brand_colors["accent_green"])
        if "text" in brand_colors:
            self._text_rgb = self._hex_to_rgb(brand_colors["text"])
        if "bg_lite" in brand_colors:
            self._zebra_rgb = self._hex_to_rgb(brand_colors["bg_lite"])

        try:
            # Si el archivo de plantilla existe, cargarlo. Si no, iniciar una en blanco.
            if os.path.exists(template_path):
                prs = Presentation(template_path)
                self._clear_slides(prs)
            else:
                prs = Presentation()
                # Widescreen 16:9 por defecto
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)

            ctx = self._engine.build_pptx_context(content)

            # Generar diapositivas a partir del contexto
            for slide_idx, slide_data in enumerate(ctx["slides"]):
                slide_type = slide_data.get("type")

                if slide_type == "cover":
                    self._add_cover_slide(prs, slide_data, brand_cfg, content)
                elif slide_type == "index":
                    self._add_index_slide(prs, slide_data)
                elif slide_type == "section":
                    self._add_section_slide(prs, slide_data, brand_cfg)
                elif slide_type == "content":
                    self._add_content_slide(prs, slide_data, brand_cfg, slide_idx)
                elif slide_type == "table":
                    self._add_table_slide(prs, slide_data, brand_cfg, slide_idx)
                elif slide_type == "closing":
                    self._add_closing_slide(prs, brand_cfg)

            # Guardar la presentación generada
            buffer = BytesIO()
            prs.save(buffer)
            logger.info(f"PPTX generado exitosamente programático: '{content.title}'")
            return buffer.getvalue()

        except Exception as e:
            logger.error(f"Error generando PPTX '{content.title}': {e}")
            raise RuntimeError(f"Error al generar el archivo PPTX: {e}") from e

    def _get_blank_layout(self, prs: Presentation):
        """Busca un slide layout completamente en blanco."""
        for layout in prs.slide_layouts:
            if len(layout.placeholders) == 0:
                return layout
        return prs.slide_layouts[-1]

    def _clear_slides(self, prs: Presentation) -> None:
        """Elimina todos los slides existentes en una presentación."""
        for i in range(len(prs.slides) - 1, -1, -1):
            slide_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[i]

    def _hex_to_rgb(self, hex_str: str) -> RGBColor:
        """Convierte un color en formato hexadecimal (#ffffff o ffffff) a RGBColor."""
        hex_str = hex_str.lstrip('#')
        return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))

    def _add_glow(self, slide, left, top, width, height, color_hex: str, opacity_percent: float) -> None:
        """Inserta un círculo translúcido (glow radial) para simular el degradado corporativo."""
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(color_hex)
        shape.line.fill.background()  # Sin borde

        # Manipulación de XML para aplicar transparencia (alpha)
        solidFill = shape.fill._fill._solidFill
        color_el = solidFill.srgbClr
        if color_el is not None:
            # Eliminar elementos alpha existentes
            for child in list(color_el):
                if child.tag.endswith('alpha'):
                    color_el.remove(child)
            # Agregar el nodo de transparencia
            alpha = parse_xml(
                f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                f'val="{int(opacity_percent * 1000)}"/>'
            )
            color_el.append(alpha)

    def _set_dark_background(self, slide) -> None:
        """Aplica un degradado diagonal oscuro premium en el slide master."""
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = 135  # Degradado diagonal
        
        # Stop 1 (Teal muy oscuro)
        stop1 = fill.gradient_stops[0]
        stop1.position = 0.0
        stop1.color.rgb = self._hex_to_rgb("#011E23")
        
        # Stop 2 (Verde bosque oscuro)
        stop2 = fill.gradient_stops[1]
        stop2.position = 1.0
        stop2.color.rgb = self._hex_to_rgb("#10473F")

    def _set_slide_background(self, slide, hex_color: str) -> None:
        """Cambia el color de fondo de una diapositiva."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(hex_color)

    def _add_header_dark_slide(self, slide, brand_cfg: dict) -> None:
        """Dibuja el encabezado premium en slides oscuros."""
        # Texto izquierdo: "Straight to Growth"
        add_textbox_text(
            slide, Inches(1.0), Inches(0.4), Inches(2.0), Inches(0.4),
            "Straight to Growth", 11, RGBColor(113, 128, 150), font_name="Georgia", italic=True
        )
        
        # Línea de conexión al lado del texto
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(0.58), Inches(7.5), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(113, 128, 150)
        line.line.color.rgb = RGBColor(113, 128, 150)

        # Logotipo blanco a la derecha
        logo_white_path = brand_cfg["logos"].get("white")
        if logo_white_path and os.path.exists(logo_white_path):
            slide.shapes.add_picture(logo_white_path, Inches(10.33), Inches(0.4), height=Inches(0.35))

    def _add_slide_footer(self, slide, slide_idx: int, prs: Presentation) -> None:
        """Inserta una línea inferior decorativa y pie de página en diapositivas interiores."""
        # Línea de separación inferior
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.8), Inches(11.33), Inches(0.01))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(217, 217, 217)  # #D9D9D9
        shape.line.color.rgb = RGBColor(217, 217, 217)

        # Texto de copyright
        add_textbox_text(
            slide, Inches(1.0), Inches(6.9), Inches(8.0), Inches(0.3),
            "© Intelligence Customer Acquisition — Convertia — Documentación interna",
            8.5, RGBColor(113, 128, 150)
        )

        # Número de diapositiva
        add_textbox_text(
            slide, Inches(10.33), Inches(6.9), Inches(2.0), Inches(0.3),
            f"Pág. {slide_idx + 1}", 8.5, RGBColor(113, 128, 150), align=PP_ALIGN.RIGHT
        )

    def _add_cover_slide(self, prs: Presentation, data: dict, brand_cfg: dict, content: DocumentContent) -> None:
        """Genera el slide de portada con el degradado oscuro corporativo, logos y texto partido."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_dark_background(slide)

        # Agregar destellos (glow) corporativos en el fondo
        self._add_glow(slide, Inches(8.0), Inches(2.5), Inches(7.5), Inches(7.5), "#8f8cff", 12)  # Glow morado
        self._add_glow(slide, Inches(-2.0), Inches(1.5), Inches(6.5), Inches(6.5), "#1AEB9F", 10)  # Glow verde menta

        # Encabezado premium
        self._add_header_dark_slide(slide, brand_cfg)

        # Separar el título dinámicamente en dos líneas para la portada premium
        title = data.get("title", content.title)
        words = title.split()
        if len(words) > 1:
            line1 = " ".join(words[:-1])
            line2 = words[-1]
        else:
            line1 = title
            line2 = ""

        # Escribir Línea 1 (Blanco) y Línea 2 (Gris claro/Lavanda)
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        # Párrafo 1
        p1 = tf.paragraphs[0]
        p1.text = line1
        p1.font.name = "Arial"
        p1.font.size = Pt(46)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.space_after = Pt(8)

        # Párrafo 2
        p2 = tf.add_paragraph()
        p2.text = line2
        p2.font.name = "Arial"
        p2.font.size = Pt(46)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(113, 128, 150)

        # Línea decorativa inferior
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.7), Inches(11.33), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(113, 128, 150)
        line.line.color.rgb = RGBColor(113, 128, 150)

    def _add_index_slide(self, prs: Presentation, data: dict) -> None:
        """Genera la diapositiva de índice/contenido con fondo blanco."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_slide_background(slide, "#FFFFFF")

        # Brillo verde sutil
        self._add_glow(slide, Inches(-2.0), Inches(3.5), Inches(5.5), Inches(5.5), "#1AEB9F", 8)

        # Columna Izquierda: Título "Contenido"
        add_textbox_text(
            slide, Inches(1.0), Inches(1.8), Inches(3.5), Inches(1.0),
            "Contenido", 36, self._primary_rgb, bold=True
        )

        # Botón redondeado con flecha ->
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.9), Inches(1.0), Inches(0.45))
        btn.fill.solid()
        btn.fill.fore_color.rgb = RGBColor(255, 255, 255)
        btn.line.color.rgb = self._primary_rgb
        btn.line.width = Pt(1.5)

        tf_btn = btn.text_frame
        tf_btn.word_wrap = False
        tf_btn.margin_left = tf_btn.margin_right = tf_btn.margin_top = tf_btn.margin_bottom = Inches(0)
        p_btn = tf_btn.paragraphs[0]
        p_btn.text = "→"
        p_btn.alignment = PP_ALIGN.CENTER
        p_btn.font.name = "Arial"
        p_btn.font.size = Pt(20)
        p_btn.font.bold = True
        p_btn.font.color.rgb = self._primary_rgb

        # Separador vertical
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.5), Inches(0.01), Inches(4.5))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(217, 217, 217)
        line.line.color.rgb = RGBColor(217, 217, 217)

        # Columna Derecha: Índice
        txBox = slide.shapes.add_textbox(Inches(5.4), Inches(1.8), Inches(6.8), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

        first = True
        for item in data.get("items", []):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"{item['number']}. {item['title']}"
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._primary_rgb
            p.space_before = Pt(8)
            p.space_after = Pt(4)

            for sub in item.get("subsections", []):
                p_sub = tf.add_paragraph()
                p_sub.text = f"    {sub['number']}.  {sub['title']}"
                p_sub.font.name = "Arial"
                p_sub.font.size = Pt(12)
                p_sub.font.color.rgb = self._text_rgb
                p_sub.space_after = Pt(3)

    def _add_section_slide(self, prs: Presentation, data: dict, brand_cfg: dict) -> None:
        """Genera una diapositiva divisora de sección (Nivel 1) con fondo oscuro, número grande y título."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_dark_background(slide)

        # Brillo verde menta a la derecha
        self._add_glow(slide, Inches(8.33), Inches(2.0), Inches(6.5), Inches(6.5), "#1AEB9F", 12)

        # Encabezado
        self._add_header_dark_slide(slide, brand_cfg)

        # Número de sección grande verde menta (en el lado izquierdo)
        num = f"{data.get('number', 1):02d}"
        add_textbox_text(slide, Inches(1.0), Inches(2.4), Inches(11.33), Inches(1.0), num, 68, self._accent_rgb, bold=True)

        # Título de sección en blanco
        title = data.get("title", "")
        add_textbox_text(slide, Inches(1.0), Inches(3.6), Inches(11.33), Inches(1.8), title, 36, RGBColor(255, 255, 255), bold=True)

        # Línea de acento verde menta abajo del título
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.4), Inches(1.5), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = self._accent_rgb
        line.line.color.rgb = self._accent_rgb

    def _add_content_slide(self, prs: Presentation, data: dict, brand_cfg: dict, slide_idx: int) -> None:
        """Genera una diapositiva de contenido con diseño de 2 columnas (Título Izquierda, Contenido Derecha)."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_slide_background(slide, "#FFFFFF")

        # Agregar destello verde menta sutil en el fondo a la izquierda
        self._add_glow(slide, Inches(-2.0), Inches(3.5), Inches(5.5), Inches(5.5), "#1AEB9F", 8)

        # -------------------------------------------------------------
        # Columna Izquierda (Título)
        # -------------------------------------------------------------
        title = data.get("title", "")
        add_textbox_text(
            slide, Inches(1.0), Inches(1.8), Inches(4.0), Inches(4.5),
            title, 34, self._primary_rgb, bold=True
        )

        # -------------------------------------------------------------
        # Columna Derecha (Cuerpo de Texto y Viñetas)
        # -------------------------------------------------------------
        txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(6.83), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)

        first = True

        # Párrafos de texto
        content_text = data.get("content", "")
        if content_text:
            for para in content_text.split('\n'):
                if para.strip():
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = para.strip()
                    p.font.name = "Arial"
                    p.font.size = Pt(14)
                    p.font.color.rgb = self._text_rgb
                    p.space_after = Pt(12)
                    p.line_spacing = 1.15

        # Viñetas principales
        bullets = data.get("bullets", [])
        if bullets:
            for bullet in bullets:
                if bullet.strip():
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = f"•  {bullet.strip()}"
                    p.font.name = "Arial"
                    p.font.size = Pt(14)
                    p.font.color.rgb = self._text_rgb
                    p.space_after = Pt(8)

        # Sub-encabezados (h3) y sus viñetas
        subheadings = data.get("subheadings", [])
        for sub in subheadings:
            p_sub = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p_sub.text = sub.get("title", "")
            p_sub.font.name = "Arial"
            p_sub.font.size = Pt(16)
            p_sub.font.bold = True
            p_sub.font.color.rgb = self._secondary_rgb
            p_sub.space_before = Pt(14)
            p_sub.space_after = Pt(6)

            sub_content = sub.get("content", "")
            if sub_content:
                for para in sub_content.split('\n'):
                    if para.strip():
                        p = tf.add_paragraph()
                        p.text = para.strip()
                        p.font.name = "Arial"
                        p.font.size = Pt(14)
                        p.font.color.rgb = self._text_rgb
                        p.space_after = Pt(8)

            sub_bullets = sub.get("bullets", [])
            for sb in sub_bullets:
                if sb.strip():
                    p = tf.add_paragraph()
                    p.text = f"•  {sb.strip()}"
                    p.font.name = "Arial"
                    p.font.size = Pt(14)
                    p.font.color.rgb = self._text_rgb
                    p.space_after = Pt(6)

        # Si hay una tabla inline dentro de la sección, insertarla abajo del contenido en la derecha
        if "table" in data:
            self._insert_table(slide, data["table"], top=Inches(3.8), height=Inches(2.5), left=Inches(5.5), width=Inches(6.83))

    def _add_table_slide(self, prs: Presentation, data: dict, brand_cfg: dict, slide_idx: int) -> None:
        """Genera una diapositiva dedicada a una tabla de datos (a pantalla completa)."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_slide_background(slide, "#FFFFFF")

        # Agregar destello verde menta sutil en el fondo a la izquierda
        self._add_glow(slide, Inches(-2.0), Inches(3.5), Inches(5.5), Inches(5.5), "#1AEB9F", 8)

        # Título
        title = data.get("title", "Datos")
        add_textbox_text(slide, Inches(1.0), Inches(1.2), Inches(11.33), Inches(0.8), title, 26, self._primary_rgb, bold=True)

        # Insertar tabla a pantalla completa
        table_data = data.get("table", {})
        if table_data and table_data.get("rows"):
            self._insert_table(slide, table_data, top=Inches(2.2), height=Inches(4.2), left=Inches(1.0), width=Inches(11.33))

    def _add_closing_slide(self, prs: Presentation, brand_cfg: dict) -> None:
        """Genera la diapositiva de cierre con fondo oscuro."""
        slide = prs.slides.add_slide(self._get_blank_layout(prs))
        self._set_dark_background(slide)

        # Agregar destellos (glow) corporativos en el fondo
        self._add_glow(slide, Inches(8.0), Inches(2.5), Inches(7.5), Inches(7.5), "#8f8cff", 12)  # Glow morado
        self._add_glow(slide, Inches(-2.0), Inches(1.5), Inches(6.5), Inches(6.5), "#1AEB9F", 10)  # Glow verde menta

        # Encabezado
        self._add_header_dark_slide(slide, brand_cfg)

        # Logo blanco centrado
        logo_white_path = brand_cfg["logos"].get("white")
        if logo_white_path and os.path.exists(logo_white_path):
            slide.shapes.add_picture(logo_white_path, Inches(4.66), Inches(2.5), width=Inches(4.0))

        # Línea decorativa inferior
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.6), Inches(11.33), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(113, 128, 150)
        line.line.color.rgb = RGBColor(113, 128, 150)

        # Textos de pie
        add_textbox_text(slide, Inches(1.0), Inches(6.8), Inches(4.0), Inches(0.4), "Straight to growth", 11, RGBColor(113, 128, 150), font_name="Georgia", italic=True)
        add_textbox_text(slide, Inches(8.33), Inches(6.8), Inches(4.0), Inches(0.4), "www.convertia.com", 11, RGBColor(113, 128, 150), align=PP_ALIGN.RIGHT)

    def _insert_table(self, slide, table_data: dict, top, height, left, width) -> None:
        """Helper para construir y aplicar estilos corporativos a una tabla en PowerPoint."""
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers:
            return

        n_cols = len(headers)
        n_rows = len(rows) + 1  # Incluye cabecera

        # Crear tabla en PowerPoint
        tshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = tshape.table

        # Estilizar Cabecera (#011E23, texto blanco y en negrita)
        for c_idx, header in enumerate(headers):
            cell = table.cell(0, c_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._primary_rgb
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)

            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.name = "Arial"
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(11)

        # Estilizar Filas de Datos con efecto Cebra
        for r_idx, row in enumerate(rows):
            bg_rgb = self._zebra_rgb if r_idx % 2 == 1 else RGBColor(255, 255, 255)

            for c_idx, val in enumerate(row):
                if c_idx >= n_cols:
                    continue
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_rgb
                cell.margin_left = Inches(0.12)
                cell.margin_right = Inches(0.12)

                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.name = "Arial"
                        run.font.color.rgb = self._text_rgb
                        run.font.bold = False
                        run.font.size = Pt(10)
